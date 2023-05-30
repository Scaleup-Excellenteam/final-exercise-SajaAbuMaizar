import json
import os
import argparse
from pptx import Presentation
import openai
import asyncio

# Set up OpenAI API credentials
openai.api_key = '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'

# Set rate limit and token count and wait time variables
RATE_LIMIT = 3  # Number of requests per minute
TOKEN_COUNT = 0
WAIT_TIME = 60


# Function to process a single slide asynchronously
async def process_slide(slide_text):
    global TOKEN_COUNT

    # Create the GPT prompt
    prompt = f"Explain the following slide:\n\n{slide_text}"

    try:
        # Check rate limit
        if TOKEN_COUNT >= RATE_LIMIT:
            wait_time = WAIT_TIME  # Wait a minute to send another request
            print(f"Rate limit exceeded. Waiting for {wait_time} seconds...")
            await asyncio.sleep(wait_time)
            TOKEN_COUNT = 0

        # Send the prompt to the OpenAI API in a separate thread and await its result
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        print(response)

        # Update token count
        TOKEN_COUNT += 1

        # Extract the AI's reply from the response
        explanation = response.choices[0].message.content

        return explanation
    except openai.error.APIError as e:
        if 'quota' in str(e):
            return "Error processing slide: Exceeded quota. Please check your plan and billing details."
        else:
            return f"Error processing slide: {str(e)}"
    except Exception as e:
        # Handle other exceptions gracefully
        return f"Error processing slide: {str(e)}"


# Function to process all slides in the presentation
async def process_presentation(presentation_path):
    # Load the presentation
    presentation = Presentation(presentation_path)

    explanations = []

    # Create a list to hold the slide processing tasks
    slide_tasks = []

    for slide in presentation.slides:
        slide_text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text

        slide_text = slide_text.strip()

        if slide_text:
            # Process the slide asynchronously and append the task to the list
            slide_task = process_slide(slide_text)
            slide_tasks.append(slide_task)

    # Wait for all slide processing tasks to complete
    explanations = await asyncio.gather(*slide_tasks)

    return explanations


# Main function
async def main():
    # Parse command-line arguments
    parser = argparse.ArgumentParser(description='GPT-Explainer')
    parser.add_argument('presentation', type=str, help='Path to the presentation file')
    args = parser.parse_args()

    presentation_path = args.presentation

    # Process the presentation
    explanations = await process_presentation(presentation_path)

    # Save the explanations to a JSON file
    output_file = os.path.splitext(presentation_path)[0] + '.json'
    with open(output_file, 'w') as f:
        json.dump(explanations, f, indent=4)

    print(f"Explanations saved to {output_file}")


# Run the program
if __name__ == '__main__':
    asyncio.run(main())
