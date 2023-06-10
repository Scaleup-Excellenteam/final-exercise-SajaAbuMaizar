import json
import os
from pptx import Presentation
import asyncio
from slide_processor import SlideProcessor

OUTPUTS_FOLDER = "outputs"


class PresentationProcessor:
    def __init__(self):
        self.slide_processor = SlideProcessor()

    async def process_presentation(self, presentation_path):
        presentation = Presentation(presentation_path)
        explanations = []

        # Create a list to hold the slide processing tasks
        slide_tasks = []

        for slide in presentation.slides:
            slide_text = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text += run.text

            slide_text = slide_text.strip()

            if slide_text:
                # Process the slide asynchronously and append the task to the list
                slide_task = self.slide_processor.process_slide(slide_text)
                slide_tasks.append(slide_task)

        # Wait for all slide processing tasks to complete
        explanations = await asyncio.gather(*slide_tasks)

        return explanations

    async def main(self, presentation_path):
        # Process the presentation
        explanations = await self.process_presentation(presentation_path)

        # Extract the filename without extension
        filename_without_extension = os.path.splitext(os.path.basename(presentation_path))[0]

        # Create the output file path in the outputs folder
        output_file_path = os.path.join(OUTPUTS_FOLDER, f"{filename_without_extension}.json")

        # Save the explanations to a JSON file in the outputs folder
        with open(output_file_path, 'w') as f:
            json.dump(explanations, f, indent=4)

        print(f"Explanations saved to {output_file_path}")
